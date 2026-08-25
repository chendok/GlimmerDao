"""文档解析服务 — 支持 PDF、Word、Markdown、TXT、Excel、ePub、PPT、HTML、图片 等格式"""
import logging
import os
from pathlib import Path
from typing import Tuple, List, Optional

logger = logging.getLogger("uvicorn")


async def parse_document(file_path: str, file_type: str, slides_dir: Optional[str] = None) -> Tuple[str, str, List[dict]]:
    """
    解析文档，返回 (纯文本内容, Markdown内容, 分页列表)
    分页列表每项: {"page_number": int, "content_text": str, "image_path": Optional[str]}
    slides_dir: PPT 幻灯片图片输出目录（仅 pptx 类型使用）
    """
    ext = file_type.lower()
    if ext == "pdf":
        return await _parse_pdf(file_path)
    elif ext == "docx":
        return await _parse_docx(file_path)
    elif ext == "md":
        return await _parse_markdown(file_path)
    elif ext == "txt":
        return await _parse_txt(file_path)
    elif ext == "xlsx":
        return await _parse_xlsx(file_path)
    elif ext == "epub":
        return await _parse_epub(file_path)
    elif ext == "mobi":
        return await _parse_mobi(file_path)
    elif ext == "pptx":
        return await _parse_pptx(file_path, slides_dir)
    elif ext == "html":
        return await _parse_html(file_path)
    elif ext in ("jpg", "jpeg", "png", "gif", "bmp", "webp", "image"):
        return await _parse_image(file_path)
    else:
        raise ValueError(f"不支持的文件格式: {file_type}")


def _get_file_type_from_extension(file_path: str) -> str:
    """根据扩展名判断文件类型"""
    ext = Path(file_path).suffix.lower().lstrip(".")
    type_map = {
        "pdf": "pdf",
        "doc": "docx",
        "docx": "docx",
        "md": "md",
        "markdown": "md",
        "txt": "txt",
        "xls": "xlsx",
        "xlsx": "xlsx",
        "epub": "epub",
        "mobi": "mobi",
        "ppt": "pptx",
        "pptx": "pptx",
        "html": "html",
        "htm": "html",
        "jpg": "image",
        "jpeg": "image",
        "png": "image",
        "gif": "image",
        "bmp": "image",
        "webp": "image",
    }
    return type_map.get(ext, ext)


# ── PDF 解析 ──

async def _parse_pdf(file_path: str) -> Tuple[str, List[dict]]:
    """使用 pdfplumber 解析 PDF"""
    try:
        import pdfplumber
    except ImportError:
        logger.warning("pdfplumber 未安装，使用简易文本提取")
        return _fallback_parse(file_path)

    pages = []
    all_text = []
    all_md = []

    try:
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ""
                pages.append({
                    "page_number": i,
                    "content_text": text,
                })
                all_text.append(text)
                all_md.append(text)

        content_text = "\n\n".join(all_text)
        content_md = "\n\n".join(all_md)
        logger.info(f"PDF 解析完成: {file_path}, 共 {len(pages)} 页")
        return content_text, content_md, pages

    except Exception as e:
        logger.error(f"PDF 解析失败: {file_path}, 错误: {e}")
        return _fallback_parse(file_path)


# ── Word 解析 ──

async def _parse_docx(file_path: str) -> Tuple[str, List[dict]]:
    """使用 python-docx 解析 Word 文档"""
    try:
        from docx import Document
    except ImportError:
        logger.warning("python-docx 未安装")
        return _fallback_parse(file_path)

    try:
        doc = Document(file_path)
        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # 根据样式判断标题层级
                if para.style.name.startswith("Heading"):
                    level = int(para.style.name.split()[-1]) if para.style.name.split()[-1].isdigit() else 1
                    prefix = "#" * min(level, 6)
                    paragraphs.append(f"{prefix} {text}")
                else:
                    paragraphs.append(text)

        # 处理表格
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
            if rows:
                # 添加表头分隔线
                header_sep = "|" + "|".join([" --- " for _ in range(len(table.rows[0].cells))]) + "|"
                rows.insert(1, header_sep)
                paragraphs.append("\n".join(rows))

        content_text = "\n\n".join(p for p in paragraphs if not p.startswith("|"))
        content_md = "\n\n".join(paragraphs)
        pages = [{"page_number": 1, "content_text": content_text}]

        logger.info(f"Word 解析完成: {file_path}")
        return content_text, content_md, pages

    except Exception as e:
        logger.error(f"Word 解析失败: {file_path}, 错误: {e}")
        return _fallback_parse(file_path)


# ── Markdown 解析 ──

async def _parse_markdown(file_path: str) -> Tuple[str, List[dict]]:
    """解析 Markdown 文件"""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            content_md = f.read()

        # 提取纯文本（去除 Markdown 标记）
        import re
        content_text = re.sub(r"#{1,6}\s+", "", content_md)  # 去除标题标记
        content_text = re.sub(r"\*\*(.+?)\*\*", r"\1", content_text)  # 粗体
        content_text = re.sub(r"\*(.+?)\*", r"\1", content_text)  # 斜体
        content_text = re.sub(r"\[(.+?)\]\(.+?\)", r"\1", content_text)  # 链接
        content_text = re.sub(r"`{1,3}[^`]*`{1,3}", "", content_text)  # 代码
        content_text = re.sub(r"\|.*\|", "", content_text)  # 表格
        content_text = re.sub(r"[-*_]{3,}", "", content_text)  # 分隔线

        pages = [{"page_number": 1, "content_text": content_text}]
        logger.info(f"Markdown 解析完成: {file_path}")
        return content_text, content_md, pages

    except Exception as e:
        logger.error(f"Markdown 解析失败: {file_path}, 错误: {e}")
        return _fallback_parse(file_path)


# ── TXT 解析 ──

async def _parse_txt(file_path: str) -> Tuple[str, List[dict]]:
    """解析纯文本文件"""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()

        content_md = content  # 纯文本即 Markdown
        pages = [{"page_number": 1, "content_text": content}]
        logger.info(f"TXT 解析完成: {file_path}")
        return content, content_md, pages

    except Exception as e:
        logger.error(f"TXT 解析失败: {file_path}, 错误: {e}")
        return _fallback_parse(file_path)


# ── Excel 解析 ──

async def _parse_xlsx(file_path: str) -> Tuple[str, List[dict]]:
    """使用 openpyxl 解析 Excel"""
    try:
        import openpyxl
    except ImportError:
        logger.warning("openpyxl 未安装")
        return _fallback_parse(file_path)

    try:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        all_sheets = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines = [f"## {sheet_name}\n"]
            rows = list(ws.iter_rows(values_only=True))

            if not rows:
                continue

            for i, row in enumerate(rows):
                cells = [str(cell) if cell is not None else "" for cell in row]
                lines.append("| " + " | ".join(cells) + " |")
                if i == 0:
                    header_sep = "|" + "|".join([" --- " for _ in range(len(cells))]) + "|"
                    lines.append(header_sep)

            all_sheets.append("\n".join(lines))

        wb.close()
        content_md = "\n\n".join(all_sheets)
        # 纯文本：去除表格标记
        import re
        content_text = re.sub(r"\|.*\|", "", content_md)
        content_text = re.sub(r"#+\s+", "", content_text)

        pages = [{"page_number": 1, "content_text": content_text}]
        logger.info(f"Excel 解析完成: {file_path}")
        return content_text, content_md, pages

    except Exception as e:
        logger.error(f"Excel 解析失败: {file_path}, 错误: {e}")
        return _fallback_parse(file_path)


# ── ePub 解析 ──

async def _parse_epub(file_path: str) -> Tuple[str, List[dict]]:
    """使用 ebooklib 解析 ePub"""
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
    except ImportError:
        logger.warning("ebooklib 或 beautifulsoup4 未安装")
        return _fallback_parse(file_path)

    try:
        book = epub.read_epub(file_path)
        chapters = []

        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_content(), "html.parser")
            text = soup.get_text(separator="\n", strip=True)
            if text:
                chapters.append(text)

        content_text = "\n\n".join(chapters)
        content_md = content_text
        pages = [{"page_number": 1, "content_text": content_text}]
        logger.info(f"ePub 解析完成: {file_path}")
        return content_text, content_md, pages

    except Exception as e:
        logger.error(f"ePub 解析失败: {file_path}, 错误: {e}")
        return _fallback_parse(file_path)


# ── Mobi 解析 ──

async def _parse_mobi(file_path: str) -> Tuple[str, List[dict]]:
    """Mobi 格式（当前回退到简易文本提取）"""
    logger.warning(f"Mobi 格式暂不支持完整解析，使用简易提取: {file_path}")
    return _fallback_parse(file_path)


# ── 图片解析 ──

async def _parse_image(file_path: str) -> Tuple[str, str, List[dict]]:
    """图片文件 — 提取元数据（原图通过 /documents/{id}/file 端点展示）"""
    try:
        from PIL import Image
        img = Image.open(file_path)
        info = f"[图片文件] 尺寸: {img.size[0]}x{img.size[1]}, 格式: {img.format}"
        pages = [{"page_number": 1, "content_text": info}]
        return info, info, pages
    except Exception as e:
        logger.error(f"图片解析失败: {file_path}, 错误: {e}")
        return _fallback_parse(file_path)


# ── PowerPoint 解析 ──

def _table_to_markdown(rows_data: List[List[str]]) -> str:
    """将二维列表转为 markdown 表格"""
    if not rows_data:
        return ""
    lines = []
    for i, row in enumerate(rows_data):
        cells = [str(c).replace("|", "\\|").replace("\n", " ").strip() for c in row]
        lines.append("| " + " | ".join(cells) + " |")
        if i == 0:
            sep = "|" + "|".join([" --- " for _ in range(len(cells))]) + "|"
            lines.append(sep)
    return "\n".join(lines)


# ── 幻灯片渲染为图片 ──

# 字体缓存
_FONT_CACHE = {}
_FONT_PATHS = [
    # Windows 中文字体
    "C:/Windows/Fonts/msyh.ttc",       # 微软雅黑
    "C:/Windows/Fonts/simhei.ttf",     # 黑体
    "C:/Windows/Fonts/simsun.ttc",     # 宋体
    "C:/Windows/Fonts/arial.ttf",      # Arial
    "C:/Windows/Fonts/calibri.ttf",    # Calibri
    # Linux 字体
    "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    # macOS 字体
    "/System/Library/Fonts/PingFang.ttc",
    "/System/Library/Fonts/Helvetica.ttc",
]


def _get_font(size: int, bold: bool = False) -> "ImageFont.FreeTypeFont":
    """获取字体，支持缓存"""
    from PIL import ImageFont
    cache_key = (size, bold)
    if cache_key in _FONT_CACHE:
        return _FONT_CACHE[cache_key]

    font = None
    for font_path in _FONT_PATHS:
        if os.path.exists(font_path):
            try:
                font = ImageFont.truetype(font_path, size)
                break
            except Exception:
                continue

    if font is None:
        # 回退到 Pillow 默认字体
        font = ImageFont.load_default()

    _FONT_CACHE[cache_key] = font
    return font


def _render_slide_to_image(slide, output_path: str, slide_width_emu: int, slide_height_emu: int, target_width: int = 1280) -> str:
    """将单张幻灯片渲染为 PNG 图片"""
    from PIL import Image, ImageDraw, ImageFont
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    import textwrap
    import io

    # 幻灯片尺寸 (EMU → 像素)
    slide_width_px = int(slide_width_emu / 9525)
    slide_height_px = int(slide_height_emu / 9525)

    scale = target_width / slide_width_px
    img_width = target_width
    img_height = int(slide_height_px * scale)

    # 创建画布
    img = Image.new("RGB", (img_width, img_height), "white")
    draw = ImageDraw.Draw(img)

    def emu2px(emu):
        return int(emu / 9525 * scale)

    def get_color(rgb_color):
        """将 pptx 颜色转为 RGB 元组"""
        try:
            return tuple(int(rgb_color[i:i+2], 16) for i in range(0, 6, 2))
        except Exception:
            return None

    def _safe_get_font_color(run):
        """安全获取 run 的字体颜色，返回 RGB 元组或 None"""
        try:
            if run.font.color and run.font.color.rgb:
                return get_color(str(run.font.color.rgb))
        except Exception:
            pass
        return None

    def draw_textbox(shape, left, top, width, height):
        """渲染文本框"""
        tf = shape.text_frame
        # 文本框背景
        fill = shape.fill
        bg_color = None
        try:
            if fill.type is not None:
                try:
                    bg_color = get_color(str(fill.fore_color.rgb))
                except Exception:
                    pass
        except Exception:
            pass

        if bg_color:
            draw.rectangle([left, top, left + width, top + height], fill=bg_color)

        # 计算文本区域（留边距）
        margin_left = emu2px(tf.margin_left or 91440)  # 默认 0.1 inch
        margin_top = emu2px(tf.margin_top or 45720)     # 默认 0.05 inch
        margin_right = emu2px(tf.margin_right or 91440)
        margin_bottom = emu2px(tf.margin_bottom or 45720)

        text_left = left + margin_left
        text_top = top + margin_top
        text_width = width - margin_left - margin_right
        text_height = height - margin_top - margin_bottom

        y_offset = text_top

        for para in tf.paragraphs:
            align = para.alignment
            # 收集所有 run 的文本和格式
            if para.runs:
                # 计算段落总宽度用于对齐
                total_text = ""
                run_formats = []
                for run in para.runs:
                    text = run.text
                    if text:
                        total_text += text
                        font_size = emu2px(run.font.size) if run.font.size else emu2px(Pt(18))
                        run_formats.append({
                            "text": text,
                            "size": max(8, font_size),
                            "bold": run.font.bold or False,
                            "color": _safe_get_font_color(run),
                        })
            else:
                total_text = para.text
                run_formats = None

            if not total_text.strip():
                y_offset += emu2px(Pt(18))
                continue

            # 文字换行
            if run_formats:
                # 有格式的文本：需要逐 run 换行
                _draw_formatted_text(draw, run_formats, text_left, y_offset,
                                    text_width, text_height, align, img_width)
                # 估算行高
                max_size = max(r["size"] for r in run_formats)
                y_offset += int(max_size * 1.4)
            else:
                # 无格式文本
                font = _get_font(emu2px(Pt(18)))
                wrapped = _wrap_text(draw, total_text, font, text_width)
                for line in wrapped:
                    bbox = draw.textbbox((0, 0), line, font=font)
                    line_width = bbox[2] - bbox[0]
                    line_height = bbox[3] - bbox[1]

                    if align == PP_ALIGN.CENTER:
                        x = text_left + (text_width - line_width) // 2
                    elif align == PP_ALIGN.RIGHT:
                        x = text_left + text_width - line_width
                    else:
                        x = text_left

                    if y_offset + line_height > text_top + text_height:
                        break
                    draw.text((x, y_offset), line, fill="black", font=font)
                    y_offset += int(line_height * 1.2)

    def _draw_formatted_text(draw, run_formats, left, y_offset, max_width, max_height, align, img_width):
        """绘制带格式的文本行"""
        from pptx.enum.text import PP_ALIGN as PA
        x = left
        line_y = y_offset

        for rf in run_formats:
            font = _get_font(rf["size"], rf["bold"])
            color = rf["color"] or (0, 0, 0)
            text = rf["text"]

            bbox = draw.textbbox((0, 0), text, font=font)
            tw = bbox[2] - bbox[0]
            th = bbox[3] - bbox[1]

            if x + tw > left + max_width:
                # 换行
                x = left
                line_y += int(max(r["size"] for r in run_formats) * 1.4)

            if line_y + th > y_offset + max_height:
                break

            draw.text((x, line_y), text, fill=color, font=font)
            x += tw

    def _wrap_text(draw, text, font, max_width):
        """文字换行"""
        import textwrap
        lines = []
        for paragraph in text.split("\n"):
            # 先尝试按字符宽度精确换行
            if not paragraph:
                lines.append("")
                continue
            current_line = ""
            for char in paragraph:
                test_line = current_line + char
                bbox = draw.textbbox((0, 0), test_line, font=font)
                if bbox[2] - bbox[0] > max_width:
                    if current_line:
                        lines.append(current_line)
                    current_line = char
                else:
                    current_line = test_line
            if current_line:
                lines.append(current_line)
        return lines

    def draw_table(shape, left, top, width, height):
        """渲染表格"""
        tbl = shape.table
        rows = len(tbl.rows)
        cols = len(tbl.columns)

        col_widths = [emu2px(tbl.columns[i].width) for i in range(cols)]
        # 缩放列宽以适应 shape 宽度
        total_col_w = sum(col_widths)
        if total_col_w > 0:
            col_scale = width / total_col_w
            col_widths = [int(w * col_scale) for w in col_widths]

        row_height = height // max(rows, 1)
        font = _get_font(max(8, min(row_height - 4, 16)))

        for r in range(rows):
            row_y = top + r * row_height
            col_x = left
            for c in range(cols):
                cell = tbl.cell(r, c)
                cell_text = cell.text.strip()
                cell_w = col_widths[c]

                # 单元格背景
                draw.rectangle([col_x, row_y, col_x + cell_w, row_y + row_height],
                             outline="lightgray", fill="white")
                # 单元格文本
                if cell_text:
                    bbox = draw.textbbox((0, 0), cell_text, font=font)
                    tw = bbox[2] - bbox[0]
                    # 截断过长文本
                    if tw > cell_w - 4:
                        while tw > cell_w - 8 and len(cell_text) > 1:
                            cell_text = cell_text[:-1]
                            bbox = draw.textbbox((0, 0), cell_text + "...", font=font)
                            tw = bbox[2] - bbox[0]
                        cell_text += "..."
                    draw.text((col_x + 2, row_y + 2), cell_text, fill="black", font=font)
                col_x += cell_w

    def draw_picture(shape, left, top, width, height):
        """渲染图片"""
        try:
            image_blob = shape.image.blob
            content_type = shape.image.content_type
            img_stream = io.BytesIO(image_blob)
            pic = Image.open(img_stream)
            # 等比缩放
            try:
                pic.thumbnail((width, height), Image.Resampling.LANCZOS)
            except AttributeError:
                pic.thumbnail((width, height), Image.LANCZOS)
            # 居中放置
            px = left + (width - pic.width) // 2
            py = top + (height - pic.height) // 2
            if pic.mode == "RGBA":
                img.paste(pic, (px, py), pic)
            else:
                img.paste(pic, (px, py))
        except Exception as e:
            logger.debug(f"幻灯片图片渲染失败: {e}")

    def draw_shape(shape):
        """递归渲染形状"""
        left = emu2px(shape.left)
        top = emu2px(shape.top)
        width = emu2px(shape.width)
        height = emu2px(shape.height)

        shape_type = shape.shape_type

        # 组合形状：递归处理子形状
        if shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                draw_shape(child)
            return

        # 绘制形状背景
        try:
            fill = shape.fill
            if fill.type is not None:
                try:
                    bg = get_color(str(fill.fore_color.rgb))
                    if bg:
                        # 圆角矩形等用圆角矩形
                        if str(shape_type) in ("ROUNDED_RECTANGLE", "AUTO_SHAPE", "ROUNDED_RECTANGLE (5)"):
                            draw.rounded_rectangle([left, top, left + width, top + height],
                                                  radius=min(width, height) // 8, fill=bg)
                        else:
                            draw.rectangle([left, top, left + width, top + height], fill=bg)
                except Exception:
                    pass
        except Exception:
            pass

        # 绘制边框
        try:
            line = shape.line
            if line.fill.type is not None:
                try:
                    lc = get_color(str(line.fill.fore_color.rgb))
                    lw = max(1, int(line.width / 9525 * scale)) if line.width else 1
                    if lc:
                        draw.rectangle([left, top, left + width, top + height],
                                     outline=lc, width=lw)
                except Exception:
                    pass
        except Exception:
            pass

        # 文本框
        if shape.has_text_frame:
            draw_textbox(shape, left, top, width, height)

        # 表格
        if shape.has_table:
            draw_table(shape, left, top, width, height)

        # 图片
        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            draw_picture(shape, left, top, width, height)

    # 按 z-order 渲染所有形状
    for shape in slide.shapes:
        draw_shape(shape)

    # 保存图片
    img.save(output_path, "PNG", optimize=True)
    return output_path


async def _parse_pptx(file_path: str, slides_dir: Optional[str] = None) -> Tuple[str, str, List[dict]]:
    """
    使用 python-pptx 解析 PowerPoint，按幻灯片分页。
    如果提供 slides_dir，则同时渲染幻灯片为图片。
    """
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx 未安装")
        return _fallback_parse(file_path)

    try:
        prs = Presentation(file_path)
        pages = []
        all_text = []
        all_md = []

        for idx, slide in enumerate(prs.slides, 1):
            slide_lines = []
            slide_text_parts = []

            for shape in slide.shapes:
                # 文本框
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs).strip()
                        if not text:
                            # 兼容无 run 的占位段落
                            text = para.text.strip()
                        if text:
                            slide_lines.append(text)
                            slide_text_parts.append(text)
                # 表格
                if shape.has_table:
                    tbl = shape.table
                    rows_data = [
                        [cell.text.strip() for cell in row.cells]
                        for row in tbl.rows
                    ]
                    md_table = _table_to_markdown(rows_data)
                    if md_table:
                        slide_lines.append(md_table)
                        # 纯文本只取单元格内容
                        for row in rows_data:
                            slide_text_parts.append(" | ".join(row))

            # 备注
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_lines.append(f"> 备注：{notes_text}")
                    slide_text_parts.append(notes_text)

            page_md = f"## 第 {idx} 页\n\n" + "\n\n".join(slide_lines) if slide_lines else f"## 第 {idx} 页\n\n（空白页）"
            page_text = "\n".join(slide_text_parts) if slide_text_parts else "（空白页）"

            page_data = {"page_number": idx, "content_text": page_text}

            # 渲染幻灯片图片
            if slides_dir:
                image_filename = f"slide_{idx:03d}.png"
                image_path = os.path.join(slides_dir, image_filename)
                try:
                    _render_slide_to_image(slide, image_path, prs.slide_width, prs.slide_height)
                    page_data["image_path"] = image_path
                    logger.debug(f"幻灯片 {idx} 渲染完成: {image_path}")
                except Exception as e:
                    logger.warning(f"幻灯片 {idx} 渲染失败: {e}")

            pages.append(page_data)
            all_md.append(page_md)
            all_text.append(page_text)

        content_text = "\n\n".join(all_text)
        content_md = "\n\n".join(all_md)
        logger.info(f"PowerPoint 解析完成: {file_path}, 共 {len(pages)} 页"
                   + (f", 已渲染 {sum(1 for p in pages if 'image_path' in p)} 张幻灯片图片" if slides_dir else ""))
        return content_text, content_md, pages

    except Exception as e:
        logger.error(f"PowerPoint 解析失败: {file_path}, 错误: {e}")
        return _fallback_parse(file_path)


# ── HTML 解析 ──

async def _parse_html(file_path: str) -> Tuple[str, str, List[dict]]:
    """使用 BeautifulSoup 解析 HTML，转换为结构化 Markdown"""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        logger.warning("beautifulsoup4 未安装")
        return _fallback_parse(file_path)

    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            html_content = f.read()

        # 优先使用 lxml，回退到内置 html.parser
        try:
            soup = BeautifulSoup(html_content, "lxml")
        except Exception:
            soup = BeautifulSoup(html_content, "html.parser")

        # 移除 script/style 等非正文标签
        for tag in soup(["script", "style", "noscript", "meta", "link", "head"]):
            tag.decompose()

        md_lines = []
        text_parts = []

        def inline_md(node) -> str:
            """将节点转为行内 markdown 文本（递归处理 a/strong/em/code 等内联元素）"""
            name = getattr(node, "name", None)
            if name is None:
                # NavigableString
                return str(node)
            if name == "a":
                text = node.get_text(strip=True)
                href = node.get("href", "")
                return f"[{text}]({href})" if href and text else text
            if name in ("strong", "b"):
                inner = "".join(inline_md(c) for c in node.children)
                return f"**{inner}**" if inner.strip() else ""
            if name in ("em", "i"):
                inner = "".join(inline_md(c) for c in node.children)
                return f"*{inner}*" if inner.strip() else ""
            if name == "code":
                return f"`{node.get_text(strip=True)}`"
            if name == "br":
                return "\n"
            # 其余行内容器：拼接子节点
            return "".join(inline_md(c) for c in node.children)

        def walk(node):
            name = getattr(node, "name", None)
            if name is None:
                # NavigableString
                text = str(node).strip()
                if text:
                    md_lines.append(text)
                    text_parts.append(text)
                return

            if name in ("h1", "h2", "h3", "h4", "h5", "h6"):
                level = int(name[1])
                text = inline_md(node).strip()
                if text:
                    md_lines.append(f"{'#' * level} {text}")
                    text_parts.append(node.get_text(strip=True))
                return

            if name == "p":
                text = inline_md(node).strip()
                if text:
                    md_lines.append(text)
                    text_parts.append(node.get_text(separator=" ", strip=True))
                return

            if name == "ul":
                for li in node.find_all("li", recursive=False):
                    text = inline_md(li).strip()
                    if text:
                        md_lines.append(f"- {text}")
                        text_parts.append(node.get_text(separator=" ", strip=True))
                return

            if name == "ol":
                for i, li in enumerate(node.find_all("li", recursive=False), 1):
                    text = inline_md(li).strip()
                    if text:
                        md_lines.append(f"{i}. {text}")
                        text_parts.append(node.get_text(separator=" ", strip=True))
                return

            if name == "table":
                rows_data = []
                for tr in node.find_all("tr"):
                    cells = [tr_cell.get_text(strip=True) for tr_cell in tr.find_all(["th", "td"])]
                    if cells:
                        rows_data.append(cells)
                md_table = _table_to_markdown(rows_data)
                if md_table:
                    md_lines.append(md_table)
                    for row in rows_data:
                        text_parts.append(" | ".join(row))
                return

            if name == "a":
                text = node.get_text(strip=True)
                href = node.get("href", "")
                if text:
                    md_lines.append(f"[{text}]({href})" if href else text)
                    text_parts.append(text)
                return

            if name == "blockquote":
                text = inline_md(node).strip()
                if text:
                    md_lines.append(f"> {text}")
                    text_parts.append(node.get_text(separator=" ", strip=True))
                return

            if name in ("br", "hr"):
                md_lines.append("\n---\n" if name == "hr" else "")
                return

            # 其余容器标签递归子节点
            for child in node.children:
                walk(child)

        # 从 body 开始（若不存在则从根节点）
        root = soup.body if soup.body else soup
        for child in root.children:
            walk(child)

        content_md = "\n\n".join(line for line in md_lines if line != "")
        content_text = "\n".join(text_parts)
        if not content_text:
            content_text = soup.get_text(separator="\n", strip=True)
        pages = [{"page_number": 1, "content_text": content_text}]
        logger.info(f"HTML 解析完成: {file_path}")
        return content_text, content_md, pages

    except Exception as e:
        logger.error(f"HTML 解析失败: {file_path}, 错误: {e}")
        return _fallback_parse(file_path)


# ── 回退解析 ──

def _fallback_parse(file_path: str) -> Tuple[str, List[dict]]:
    """简易回退：尝试以文本方式读取"""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
    except (UnicodeDecodeError, Exception):
        content = f"[无法解析的文件: {os.path.basename(file_path)}]"

    pages = [{"page_number": 1, "content_text": content}]
    return content, content, pages